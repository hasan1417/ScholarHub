"""Capture each slide of the HTML deck as a 1920x1080 PNG, then assemble into a 16:9 PPTX."""
from pathlib import Path
from playwright.sync_api import sync_playwright
from pptx import Presentation
from pptx.util import Inches, Emu

HERE = Path(__file__).parent
DECK_HTML = HERE.parent / "ScholarHub Product Deck.html"
PNG_DIR = HERE / "slides"
OUT_PPTX = HERE.parent / "ScholarHub-Product-Deck.pptx"
N_SLIDES = 35

PNG_DIR.mkdir(exist_ok=True)

# --- Capture ---
print(f"Capturing {N_SLIDES} slides...")
with sync_playwright() as p:
    browser = p.chromium.launch()
    ctx = browser.new_context(viewport={"width": 1920, "height": 1080}, device_scale_factor=2)
    page = ctx.new_page()
    page.goto(f"file://{DECK_HTML}")
    page.wait_for_selector("deck-stage", timeout=10_000)
    # Wait for fonts + initial render
    page.wait_for_timeout(1500)

    # Inject CSS to hide deck-stage internal nav (the host's shadow root)
    page.evaluate("""
    () => {
      const ds = document.querySelector('deck-stage');
      const s = document.createElement('style');
      s.textContent = `
        [class*="overlay"], [class*="controls"], [class*="hud"],
        button, .nav, .progress, .indicator, .status, .toolbar { display: none !important; }
        :host { background: transparent !important; }
        .stage { padding: 0 !important; }
        .canvas { box-shadow: none !important; }
      `;
      ds.shadowRoot.appendChild(s);
    }
    """)

    for i in range(N_SLIDES):
        page.evaluate(f"document.querySelector('deck-stage').goTo({i})")
        page.wait_for_timeout(450)  # let transition + screenshot images settle
        out = PNG_DIR / f"slide-{i+1:02d}.png"
        page.screenshot(path=str(out), full_page=False, clip={"x": 0, "y": 0, "width": 1920, "height": 1080})
        print(f"  captured slide {i+1:02d}")

    browser.close()

# --- Build PPTX ---
print("Building PPTX...")
prs = Presentation()
# 16:9 — 13.333" x 7.5" at default 96 DPI
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

blank = prs.slide_layouts[6]  # blank layout
for i in range(N_SLIDES):
    img = PNG_DIR / f"slide-{i+1:02d}.png"
    if not img.exists():
        print(f"  WARN missing {img}")
        continue
    slide = prs.slides.add_slide(blank)
    slide.shapes.add_picture(str(img), 0, 0, width=prs.slide_width, height=prs.slide_height)

prs.save(str(OUT_PPTX))
size_mb = OUT_PPTX.stat().st_size / (1024 * 1024)
print(f"OK saved {OUT_PPTX.name} ({size_mb:.1f} MB, {N_SLIDES} slides)")
